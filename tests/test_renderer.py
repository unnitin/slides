"""
tests/test_renderer.py — Tests for Phase 3: PPTX Renderer + Format Plugins

Covers:
  - render() public API: file creation, slide count, filename sanitization
  - Per-type renderers: all 11 slide types produce correct shapes
  - Color utilities: resolve_color, text/muted color for backgrounds
  - Background application: dark, light, gradient
  - Text helpers: textbox creation, bullet list creation
  - Speaker notes
  - Dispatch table completeness
  - Format plugins: converter registry, PDF/EE4P converters
"""

from pathlib import Path

import pytest
from pptx import Presentation as PptxPresentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

from src.dsl.models import (
    BackgroundType,
    BrandConfig,
    BulletItem,
    ColumnContent,
    CompareTable,
    PresentationMeta,
    PresentationNode,
    SlideNode,
    SlideType,
    StatItem,
    TimelineStep,
)
from src.renderer.format_plugins import (
    EE4PConverter,
    PDFConverter,
    get_converter,
)
from src.renderer.pptx_renderer import (
    SLIDE_HEIGHT,
    SLIDE_WIDTH,
    _add_bullet_list,
    _add_speaker_notes,
    _add_textbox,
    _apply_background,
    _muted_color_for_bg,
    _RENDERERS,
    _render_bullet_points,
    _render_closing,
    _render_comparison,
    _render_freeform,
    _render_image_text,
    _render_quote,
    _render_section_divider,
    _render_stat_callout,
    _render_timeline,
    _render_title,
    _render_two_column,
    _text_color_for_bg,
    render,
    resolve_color,
)


# ── Fixtures ─────────────────────────────────────────────────────


@pytest.fixture
def brand():
    return BrandConfig()


@pytest.fixture
def blank_slide():
    """Return a blank slide from a fresh presentation."""
    prs = PptxPresentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


@pytest.fixture
def output_dir(tmp_path):
    return tmp_path


def _make_presentation(slides, **meta_kwargs):
    """Helper to build a PresentationNode."""
    meta = PresentationMeta(**meta_kwargs) if meta_kwargs else PresentationMeta()
    return PresentationNode(meta=meta, slides=slides)


# ── resolve_color ────────────────────────────────────────────────


class TestResolveColor:
    def test_resolve_primary(self, brand):
        c = resolve_color("primary", brand)
        assert c == RGBColor(0x1E, 0x27, 0x61)

    def test_resolve_secondary(self, brand):
        c = resolve_color("secondary", brand)
        assert c == RGBColor(0xCA, 0xDC, 0xFC)

    def test_resolve_accent(self, brand):
        c = resolve_color("accent", brand)
        assert c == RGBColor(0xF9, 0x61, 0x67)

    def test_resolve_named_white(self, brand):
        c = resolve_color("white", brand)
        assert c == RGBColor(0xFF, 0xFF, 0xFF)

    def test_resolve_named_black(self, brand):
        c = resolve_color("black", brand)
        assert c == RGBColor(0x00, 0x00, 0x00)

    def test_resolve_hex_direct(self, brand):
        c = resolve_color("FF8800", brand)
        assert c == RGBColor(0xFF, 0x88, 0x00)


# ── Text Color for Background ───────────────────────────────────


class TestTextColorForBg:
    def test_dark_bg_white_text(self, brand):
        c = _text_color_for_bg(BackgroundType.DARK, brand)
        assert c == RGBColor(0xFF, 0xFF, 0xFF)

    def test_gradient_bg_white_text(self, brand):
        c = _text_color_for_bg(BackgroundType.GRADIENT, brand)
        assert c == RGBColor(0xFF, 0xFF, 0xFF)

    def test_light_bg_black_text(self, brand):
        c = _text_color_for_bg(BackgroundType.LIGHT, brand)
        assert c == RGBColor(0x00, 0x00, 0x00)


class TestMutedColorForBg:
    def test_dark_bg_muted(self, brand):
        c = _muted_color_for_bg(BackgroundType.DARK, brand)
        assert c == RGBColor(0xCC, 0xCC, 0xCC)

    def test_light_bg_muted(self, brand):
        c = _muted_color_for_bg(BackgroundType.LIGHT, brand)
        assert c == RGBColor(0x66, 0x66, 0x66)


# ── Background Application ───────────────────────────────────────


class TestApplyBackground:
    def test_dark_uses_primary(self, blank_slide, brand):
        _apply_background(blank_slide, BackgroundType.DARK, brand)
        fill = blank_slide.background.fill
        assert fill.fore_color.rgb == RGBColor(0x1E, 0x27, 0x61)

    def test_light_uses_white(self, blank_slide, brand):
        _apply_background(blank_slide, BackgroundType.LIGHT, brand)
        fill = blank_slide.background.fill
        assert fill.fore_color.rgb == RGBColor(0xFF, 0xFF, 0xFF)

    def test_gradient_uses_primary(self, blank_slide, brand):
        _apply_background(blank_slide, BackgroundType.GRADIENT, brand)
        fill = blank_slide.background.fill
        assert fill.fore_color.rgb == RGBColor(0x1E, 0x27, 0x61)


# ── Text Helpers ─────────────────────────────────────────────────


class TestAddTextbox:
    def test_creates_shape(self, blank_slide):
        initial = len(blank_slide.shapes)
        _add_textbox(blank_slide, 1.0, 1.0, 5.0, 1.0, "Hello")
        assert len(blank_slide.shapes) == initial + 1

    def test_text_content(self, blank_slide):
        box = _add_textbox(blank_slide, 1.0, 1.0, 5.0, 1.0, "Test Text")
        assert box.text_frame.paragraphs[0].text == "Test Text"

    def test_bold_flag(self, blank_slide):
        box = _add_textbox(blank_slide, 1.0, 1.0, 5.0, 1.0, "Bold", bold=True)
        assert box.text_frame.paragraphs[0].font.bold is True

    def test_font_size(self, blank_slide):
        box = _add_textbox(blank_slide, 1.0, 1.0, 5.0, 1.0, "Big", font_size=36)
        assert box.text_frame.paragraphs[0].font.size == Pt(36)


class TestAddBulletList:
    def test_creates_bullets(self, blank_slide):
        bullets = [
            BulletItem(text="First"),
            BulletItem(text="Second"),
            BulletItem(text="Third"),
        ]
        box = _add_bullet_list(blank_slide, 1.0, 1.0, 5.0, 3.0, bullets)
        paragraphs = box.text_frame.paragraphs
        assert len(paragraphs) == 3
        assert paragraphs[0].text == "First"
        assert paragraphs[2].text == "Third"

    def test_sub_bullets_have_level(self, blank_slide):
        bullets = [
            BulletItem(text="Top", level=0),
            BulletItem(text="Sub", level=1),
        ]
        box = _add_bullet_list(blank_slide, 1.0, 1.0, 5.0, 2.0, bullets)
        assert box.text_frame.paragraphs[1].level == 1

    def test_sub_bullets_smaller_font(self, blank_slide):
        bullets = [
            BulletItem(text="Top", level=0),
            BulletItem(text="Sub", level=1),
        ]
        box = _add_bullet_list(
            blank_slide,
            1.0,
            1.0,
            5.0,
            2.0,
            bullets,
            font_size=14,
        )
        assert box.text_frame.paragraphs[0].font.size == Pt(14)
        assert box.text_frame.paragraphs[1].font.size == Pt(12)


# ── Dispatch Table ───────────────────────────────────────────────


class TestDispatchTable:
    def test_all_slide_types_covered(self):
        for st in SlideType:
            assert st in _RENDERERS, f"Missing renderer for {st}"

    def test_dispatch_maps_to_functions(self):
        assert _RENDERERS[SlideType.TITLE] is _render_title
        assert _RENDERERS[SlideType.STAT_CALLOUT] is _render_stat_callout
        assert _RENDERERS[SlideType.TIMELINE] is _render_timeline
        assert _RENDERERS[SlideType.COMPARISON] is _render_comparison
        assert _RENDERERS[SlideType.QUOTE] is _render_quote


# ── Per-Type Renderers ───────────────────────────────────────────


class TestRenderTitle:
    def test_adds_heading_and_subtitle(self, blank_slide, brand):
        node = SlideNode(
            slide_name="title_slide",
            slide_type=SlideType.TITLE,
            heading="Big Title",
            subheading="A subtitle here",
        )
        _render_title(blank_slide, node, brand)
        texts = [s.text_frame.paragraphs[0].text for s in blank_slide.shapes if s.has_text_frame]
        assert "Big Title" in texts
        assert "A subtitle here" in texts

    def test_no_heading_no_crash(self, blank_slide, brand):
        node = SlideNode(slide_name="empty_title", slide_type=SlideType.TITLE)
        _render_title(blank_slide, node, brand)  # should not raise


class TestRenderSectionDivider:
    def test_adds_heading_and_accent_strip(self, blank_slide, brand):
        node = SlideNode(
            slide_name="section",
            slide_type=SlideType.SECTION_DIVIDER,
            heading="Section One",
        )
        _render_section_divider(blank_slide, node, brand)
        # Heading + accent strip = 2 shapes
        assert len(blank_slide.shapes) >= 2


class TestRenderStatCallout:
    def test_renders_stats(self, blank_slide, brand):
        node = SlideNode(
            slide_name="stats",
            slide_type=SlideType.STAT_CALLOUT,
            heading="Key Metrics",
            stats=[
                StatItem(value="94%", label="Uptime", description="Up from 87%"),
                StatItem(value="3.2B", label="Events/Day"),
            ],
        )
        _render_stat_callout(blank_slide, node, brand)
        texts = [s.text_frame.paragraphs[0].text for s in blank_slide.shapes if s.has_text_frame]
        assert "94%" in texts
        assert "3.2B" in texts
        assert "Uptime" in texts

    def test_no_stats_early_return(self, blank_slide, brand):
        node = SlideNode(
            slide_name="empty_stats",
            slide_type=SlideType.STAT_CALLOUT,
            heading="No Stats Here",
        )
        initial = len(blank_slide.shapes)
        _render_stat_callout(blank_slide, node, brand)
        # Only heading textbox added
        assert len(blank_slide.shapes) == initial + 1


class TestRenderBulletPoints:
    def test_renders_bullets(self, blank_slide, brand):
        node = SlideNode(
            slide_name="bullets",
            slide_type=SlideType.BULLET_POINTS,
            heading="Key Points",
            bullets=[
                BulletItem(text="Point A"),
                BulletItem(text="Point B"),
            ],
        )
        _render_bullet_points(blank_slide, node, brand)
        texts = []
        for s in blank_slide.shapes:
            if s.has_text_frame:
                for p in s.text_frame.paragraphs:
                    texts.append(p.text)
        assert "Point A" in texts

    def test_icon_rows_layout(self, blank_slide, brand):
        node = SlideNode(
            slide_name="icon_bullets",
            slide_type=SlideType.BULLET_POINTS,
            layout="icon_rows",
            heading="Features",
            bullets=[
                BulletItem(text="Fast", icon="rocket"),
                BulletItem(text="Safe", icon="shield"),
            ],
        )
        _render_bullet_points(blank_slide, node, brand)
        # Should have heading + 2 circles + 2 texts = at least 5 shapes
        assert len(blank_slide.shapes) >= 5


class TestRenderTwoColumn:
    def test_renders_columns_with_divider(self, blank_slide, brand):
        node = SlideNode(
            slide_name="twocol",
            slide_type=SlideType.TWO_COLUMN,
            heading="Two Columns",
            columns=[
                ColumnContent(title="Left", bullets=[BulletItem(text="L1")]),
                ColumnContent(title="Right", bullets=[BulletItem(text="R1")]),
            ],
        )
        _render_two_column(blank_slide, node, brand)
        texts = []
        for s in blank_slide.shapes:
            if s.has_text_frame:
                for p in s.text_frame.paragraphs:
                    texts.append(p.text)
        assert "Left" in texts
        assert "Right" in texts
        # Divider line shape should exist
        assert len(blank_slide.shapes) >= 5  # heading + 2 titles + 2 bullets + divider


class TestRenderComparison:
    def test_renders_table(self, blank_slide, brand):
        node = SlideNode(
            slide_name="compare",
            slide_type=SlideType.COMPARISON,
            heading="Before vs After",
            compare=CompareTable(
                headers=["Feature", "Before", "After"],
                rows=[
                    ["Speed", "Slow", "Fast"],
                    ["Cost", "High", "Low"],
                ],
            ),
        )
        _render_comparison(blank_slide, node, brand)
        # Find the table shape
        table_shapes = [s for s in blank_slide.shapes if s.has_table]
        assert len(table_shapes) == 1
        table = table_shapes[0].table
        assert len(table.rows) == 3  # 1 header + 2 data
        assert len(table.columns) == 3
        assert table.cell(0, 0).text == "Feature"
        assert table.cell(1, 0).text == "Speed"

    def test_no_compare_early_return(self, blank_slide, brand):
        node = SlideNode(
            slide_name="no_compare",
            slide_type=SlideType.COMPARISON,
            heading="Empty Compare",
        )
        _render_comparison(blank_slide, node, brand)
        table_shapes = [s for s in blank_slide.shapes if s.has_table]
        assert len(table_shapes) == 0


class TestRenderTimeline:
    def test_renders_steps(self, blank_slide, brand):
        node = SlideNode(
            slide_name="timeline",
            slide_type=SlideType.TIMELINE,
            heading="Roadmap",
            timeline=[
                TimelineStep(time="Q1", title="Plan", description="Research"),
                TimelineStep(time="Q2", title="Build"),
                TimelineStep(time="Q3", title="Launch"),
            ],
        )
        _render_timeline(blank_slide, node, brand)
        texts = []
        for s in blank_slide.shapes:
            if s.has_text_frame:
                for p in s.text_frame.paragraphs:
                    texts.append(p.text)
        assert "Q1" in texts
        assert "Plan" in texts
        assert "Launch" in texts

    def test_no_timeline_early_return(self, blank_slide, brand):
        node = SlideNode(
            slide_name="no_tl",
            slide_type=SlideType.TIMELINE,
            heading="Empty Timeline",
        )
        _render_timeline(blank_slide, node, brand)
        # Only heading shape
        assert len(blank_slide.shapes) == 1


class TestRenderImageText:
    def test_missing_image_uses_placeholder(self, blank_slide, brand):
        node = SlideNode(
            slide_name="img",
            slide_type=SlideType.IMAGE_TEXT,
            heading="Architecture",
            image="/nonexistent/image.png",
            bullets=[BulletItem(text="Detail A")],
        )
        _render_image_text(blank_slide, node, brand)
        texts = []
        for s in blank_slide.shapes:
            if s.has_text_frame:
                for p in s.text_frame.paragraphs:
                    texts.append(p.text)
        assert any("[Image:" in t for t in texts)
        assert "Detail A" in texts

    def test_body_fallback(self, blank_slide, brand):
        node = SlideNode(
            slide_name="img_body",
            slide_type=SlideType.IMAGE_TEXT,
            heading="Chart",
            body="Some explanation text",
        )
        _render_image_text(blank_slide, node, brand)
        texts = [s.text_frame.paragraphs[0].text for s in blank_slide.shapes if s.has_text_frame]
        assert "Some explanation text" in texts


class TestRenderQuote:
    def test_renders_quote_and_attribution(self, blank_slide, brand):
        node = SlideNode(
            slide_name="quote",
            slide_type=SlideType.QUOTE,
            heading="Data is the new oil.",
            subheading="Clive Humby",
        )
        _render_quote(blank_slide, node, brand)
        texts = []
        for s in blank_slide.shapes:
            if s.has_text_frame:
                for p in s.text_frame.paragraphs:
                    texts.append(p.text)
        assert "Data is the new oil." in texts
        assert any("Clive Humby" in t for t in texts)
        assert any("\u201c" in t for t in texts)  # opening quote mark


class TestRenderClosing:
    def test_renders_heading_and_subtitle(self, blank_slide, brand):
        node = SlideNode(
            slide_name="closing",
            slide_type=SlideType.CLOSING,
            heading="Thank You",
            subheading="Questions?",
        )
        _render_closing(blank_slide, node, brand)
        texts = [s.text_frame.paragraphs[0].text for s in blank_slide.shapes if s.has_text_frame]
        assert "Thank You" in texts
        assert "Questions?" in texts


class TestRenderFreeform:
    def test_renders_bullets(self, blank_slide, brand):
        node = SlideNode(
            slide_name="free",
            slide_type=SlideType.FREEFORM,
            heading="Misc",
            bullets=[BulletItem(text="Anything")],
        )
        _render_freeform(blank_slide, node, brand)
        texts = []
        for s in blank_slide.shapes:
            if s.has_text_frame:
                for p in s.text_frame.paragraphs:
                    texts.append(p.text)
        assert "Anything" in texts

    def test_renders_body_when_no_bullets(self, blank_slide, brand):
        node = SlideNode(
            slide_name="free_body",
            slide_type=SlideType.FREEFORM,
            heading="Misc",
            body="Some body text",
        )
        _render_freeform(blank_slide, node, brand)
        texts = [s.text_frame.paragraphs[0].text for s in blank_slide.shapes if s.has_text_frame]
        assert "Some body text" in texts


# ── Speaker Notes ────────────────────────────────────────────────


class TestSpeakerNotes:
    def test_adds_notes(self, blank_slide):
        _add_speaker_notes(blank_slide, "Talk about X, Y, Z")
        assert blank_slide.notes_slide.notes_text_frame.text == "Talk about X, Y, Z"


# ── render() Public API ──────────────────────────────────────────


class TestRenderPublicAPI:
    def test_creates_pptx_file(self, output_dir):
        pres = _make_presentation(
            [SlideNode(slide_name="s1", slide_type=SlideType.TITLE, heading="Hello")],
            title="Test Deck",
        )
        path = render(pres, output_dir)
        assert path.exists()
        assert path.suffix == ".pptx"

    def test_correct_slide_count(self, output_dir):
        slides = [
            SlideNode(slide_name="s1", slide_type=SlideType.TITLE, heading="Title"),
            SlideNode(
                slide_name="s2",
                slide_type=SlideType.BULLET_POINTS,
                heading="Bullets",
                bullets=[BulletItem(text="A")],
            ),
            SlideNode(slide_name="s3", slide_type=SlideType.CLOSING, heading="End"),
        ]
        pres = _make_presentation(slides, title="Three Slides")
        path = render(pres, output_dir)
        opened = PptxPresentation(str(path))
        assert len(opened.slides) == 3

    def test_filename_from_title(self, output_dir):
        pres = _make_presentation(
            [SlideNode(slide_name="s1", slide_type=SlideType.TITLE, heading="Hi")],
            title="My Cool Deck",
        )
        path = render(pres, output_dir)
        assert path.name == "My Cool Deck.pptx"

    def test_filename_sanitizes_special_chars(self, output_dir):
        pres = _make_presentation(
            [SlideNode(slide_name="s1", slide_type=SlideType.TITLE, heading="Hi")],
            title="Q3/Q4 Results: A&B <Test>",
        )
        path = render(pres, output_dir)
        assert "/" not in path.name
        assert ":" not in path.name
        assert "<" not in path.name

    def test_creates_output_dir_if_missing(self, tmp_path):
        new_dir = tmp_path / "nested" / "output"
        pres = _make_presentation(
            [SlideNode(slide_name="s1", slide_type=SlideType.TITLE, heading="X")],
            title="Test",
        )
        path = render(pres, new_dir)
        assert path.exists()
        assert new_dir.exists()

    def test_slide_dimensions_16_9(self, output_dir):
        pres = _make_presentation(
            [SlideNode(slide_name="s1", slide_type=SlideType.TITLE, heading="Hi")],
            title="Wide",
        )
        path = render(pres, output_dir)
        opened = PptxPresentation(str(path))
        assert opened.slide_width == SLIDE_WIDTH
        assert opened.slide_height == SLIDE_HEIGHT

    def test_speaker_notes_rendered(self, output_dir):
        pres = _make_presentation(
            [
                SlideNode(
                    slide_name="s1",
                    slide_type=SlideType.TITLE,
                    heading="Hi",
                    speaker_notes="Remember to smile",
                )
            ],
            title="Notes Test",
        )
        path = render(pres, output_dir)
        opened = PptxPresentation(str(path))
        notes_text = opened.slides[0].notes_slide.notes_text_frame.text
        assert "Remember to smile" in notes_text

    def test_dark_background_applied(self, output_dir):
        pres = _make_presentation(
            [
                SlideNode(
                    slide_name="s1",
                    slide_type=SlideType.TITLE,
                    heading="Dark",
                    background=BackgroundType.DARK,
                )
            ],
            title="Dark BG",
        )
        path = render(pres, output_dir)
        opened = PptxPresentation(str(path))
        bg_fill = opened.slides[0].background.fill
        assert bg_fill.fore_color.rgb == RGBColor(0x1E, 0x27, 0x61)

    def test_all_slide_types_render_without_error(self, output_dir):
        """Smoke test: every slide type renders without crashing."""
        slides = [
            SlideNode(
                slide_name="title", slide_type=SlideType.TITLE, heading="Title", subheading="Sub"
            ),
            SlideNode(
                slide_name="section", slide_type=SlideType.SECTION_DIVIDER, heading="Section"
            ),
            SlideNode(
                slide_name="stats",
                slide_type=SlideType.STAT_CALLOUT,
                heading="Stats",
                stats=[StatItem(value="99%", label="Accuracy")],
            ),
            SlideNode(
                slide_name="bullets",
                slide_type=SlideType.BULLET_POINTS,
                heading="Points",
                bullets=[BulletItem(text="A"), BulletItem(text="B")],
            ),
            SlideNode(
                slide_name="twocol",
                slide_type=SlideType.TWO_COLUMN,
                heading="Columns",
                columns=[
                    ColumnContent(title="L", bullets=[BulletItem(text="L1")]),
                    ColumnContent(title="R", bullets=[BulletItem(text="R1")]),
                ],
            ),
            SlideNode(
                slide_name="compare",
                slide_type=SlideType.COMPARISON,
                heading="Compare",
                compare=CompareTable(headers=["A", "B"], rows=[["1", "2"]]),
            ),
            SlideNode(
                slide_name="timeline",
                slide_type=SlideType.TIMELINE,
                heading="Roadmap",
                timeline=[TimelineStep(time="Q1", title="Start")],
            ),
            SlideNode(
                slide_name="img",
                slide_type=SlideType.IMAGE_TEXT,
                heading="Image",
                body="Description",
            ),
            SlideNode(
                slide_name="quote",
                slide_type=SlideType.QUOTE,
                heading="A wise quote",
                subheading="Author",
            ),
            SlideNode(slide_name="closing", slide_type=SlideType.CLOSING, heading="Goodbye"),
            SlideNode(
                slide_name="free",
                slide_type=SlideType.FREEFORM,
                heading="Freeform",
                body="Anything goes",
            ),
        ]
        pres = _make_presentation(slides, title="All Types")
        path = render(pres, output_dir)
        opened = PptxPresentation(str(path))
        assert len(opened.slides) == 11

    def test_custom_brand_colors(self, output_dir):
        meta = PresentationMeta(
            title="Custom Brand",
            brand=BrandConfig(primary="FF0000", secondary="00FF00", accent="0000FF"),
        )
        pres = PresentationNode(
            meta=meta,
            slides=[
                SlideNode(
                    slide_name="s1",
                    slide_type=SlideType.TITLE,
                    heading="Red",
                    background=BackgroundType.DARK,
                )
            ],
        )
        path = render(pres, output_dir)
        opened = PptxPresentation(str(path))
        bg_fill = opened.slides[0].background.fill
        assert bg_fill.fore_color.rgb == RGBColor(0xFF, 0x00, 0x00)


# ── Format Plugins ───────────────────────────────────────────────


class TestFormatPlugins:
    def test_pdf_converter_can_convert(self):
        c = PDFConverter()
        assert c.can_convert("pdf") is True
        assert c.can_convert("PDF") is True
        assert c.can_convert("pptx") is False

    def test_ee4p_converter_can_convert(self):
        c = EE4PConverter()
        assert c.can_convert("ee4p") is True
        assert c.can_convert("EE4P") is True
        assert c.can_convert("pdf") is False

    def test_ee4p_convert_raises_not_implemented(self):
        c = EE4PConverter()
        with pytest.raises(NotImplementedError):
            c.convert(Path("test.pptx"), Path("test.ee4p"))

    def test_get_converter_pdf(self):
        c = get_converter("pdf")
        assert isinstance(c, PDFConverter)

    def test_get_converter_ee4p(self):
        c = get_converter("ee4p")
        assert isinstance(c, EE4PConverter)

    def test_get_converter_unknown_raises(self):
        with pytest.raises(ValueError, match="No converter available"):
            get_converter("docx")
