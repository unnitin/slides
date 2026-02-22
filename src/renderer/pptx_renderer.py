"""
src/renderer/pptx_renderer.py -- PPTX Rendering Engine

Converts a PresentationNode into a .pptx file using python-pptx.
Deterministic: same input always produces the same output.

Two modes:
  1. Template-based: maps slides to existing template layouts
  2. Brand-based: generates slides from scratch using BrandConfig
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from src.dsl.models import (
    BackgroundType,
    BrandConfig,
    BulletItem,
    PresentationNode,
    SlideNode,
    SlideType,
)

logger = logging.getLogger(__name__)

# ── Geometry Constants (inches, 16:9) ─────────────────────────────

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

MARGIN_TOP = 0.6
MARGIN_BOTTOM = 0.6  # Symmetric with top
MARGIN_LEFT = 0.7
MARGIN_RIGHT = 0.7

CONTENT_WIDTH = 13.333 - MARGIN_LEFT - MARGIN_RIGHT  # 11.933
CONTENT_HEIGHT = 7.5 - MARGIN_TOP - MARGIN_BOTTOM  # 6.3
CONTENT_TOP = MARGIN_TOP + 0.8  # 1.4

# Source / footnote zone at bottom of slide
SOURCE_FONT = 9
SOURCE_TOP = 7.5 - MARGIN_BOTTOM - 0.4  # ~6.5
SOURCE_HEIGHT = 0.3
PAGE_NUM_WIDTH = 0.5

TITLE_LEFT = MARGIN_LEFT
TITLE_TOP = MARGIN_TOP
TITLE_WIDTH = CONTENT_WIDTH
TITLE_HEIGHT = 0.7

# Font sizes (points)
FONT_TITLE = 36
FONT_SUBTITLE = 20
FONT_HEADING = 24
FONT_BODY = 14
FONT_CAPTION = 11
FONT_STAT_VALUE = 54
FONT_STAT_LABEL = 16
FONT_STAT_DESC = 12

# Spacing
ELEMENT_GAP = 0.3
COLUMN_GAP = 0.4


# ── Color Utilities ───────────────────────────────────────────────


def resolve_color(color_ref: str, brand: BrandConfig) -> RGBColor:
    """Resolve a color reference to an RGBColor."""
    named = {"white": "FFFFFF", "black": "000000"}
    brand_map = {
        "primary": brand.primary,
        "secondary": brand.secondary,
        "accent": brand.accent,
    }
    hex_val = brand_map.get(color_ref) or named.get(color_ref) or color_ref
    r = int(hex_val[0:2], 16)
    g = int(hex_val[2:4], 16)
    b = int(hex_val[4:6], 16)
    return RGBColor(r, g, b)


def _text_color_for_bg(bg: BackgroundType, brand: BrandConfig) -> RGBColor:
    """Return appropriate text color for a background type."""
    if bg in (BackgroundType.DARK, BackgroundType.GRADIENT):
        return RGBColor(0xFF, 0xFF, 0xFF)
    return resolve_color("black", brand)


def _muted_color_for_bg(bg: BackgroundType, brand: BrandConfig) -> RGBColor:
    """Return a muted/secondary text color for a background type."""
    if bg in (BackgroundType.DARK, BackgroundType.GRADIENT):
        return RGBColor(0xCC, 0xCC, 0xCC)
    return RGBColor(0x66, 0x66, 0x66)


# ── Background ────────────────────────────────────────────────────


def _apply_background(slide, bg_type: BackgroundType, brand: BrandConfig):
    """Apply background fill to a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()

    if bg_type == BackgroundType.DARK:
        fill.fore_color.rgb = resolve_color("primary", brand)
    elif bg_type == BackgroundType.GRADIENT:
        fill.fore_color.rgb = resolve_color("primary", brand)
    elif bg_type == BackgroundType.LIGHT:
        fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    else:
        fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


# ── Text Helpers ──────────────────────────────────────────────────


def _add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    font_size: int = FONT_BODY,
    bold: bool = False,
    color: Optional[RGBColor] = None,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
    font_name: Optional[str] = None,
) -> object:
    """Add a textbox to a slide and return the shape."""
    txBox = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.alignment = alignment
    if color:
        p.font.color.rgb = color
    if font_name:
        p.font.name = font_name
    return txBox


def _add_bullet_list(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    bullets: list[BulletItem],
    font_size: int = FONT_BODY,
    color: Optional[RGBColor] = None,
    font_name: Optional[str] = None,
) -> object:
    """Add a bulleted list to a slide."""
    txBox = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if bullet.icon:
            p.text = f"{bullet.text}"
        else:
            p.text = bullet.text

        p.font.size = Pt(font_size - (bullet.level * 2))
        p.level = bullet.level
        if color:
            p.font.color.rgb = color
        if font_name:
            p.font.name = font_name

    return txBox


# ── Per-Type Renderers ────────────────────────────────────────────


def _render_title(slide, node: SlideNode, brand: BrandConfig):
    """Render a title slide."""
    text_color = _text_color_for_bg(node.background, brand)
    muted = _muted_color_for_bg(node.background, brand)

    if node.heading:
        _add_textbox(
            slide,
            MARGIN_LEFT,
            2.0,
            CONTENT_WIDTH,
            1.5,
            node.heading,
            font_size=FONT_TITLE + 8,
            bold=True,
            color=text_color,
            alignment=PP_ALIGN.CENTER,
            font_name=brand.header_font,
        )

    if node.subheading:
        _add_textbox(
            slide,
            MARGIN_LEFT,
            3.8,
            CONTENT_WIDTH,
            1.0,
            node.subheading,
            font_size=FONT_SUBTITLE,
            color=muted,
            alignment=PP_ALIGN.CENTER,
            font_name=brand.body_font,
        )


def _render_section_divider(slide, node: SlideNode, brand: BrandConfig):
    """Render a section divider slide."""
    text_color = _text_color_for_bg(node.background, brand)

    if node.heading:
        _add_textbox(
            slide,
            MARGIN_LEFT,
            2.5,
            CONTENT_WIDTH,
            1.5,
            node.heading,
            font_size=FONT_TITLE,
            bold=True,
            color=text_color,
            alignment=PP_ALIGN.CENTER,
            font_name=brand.header_font,
        )

    # Accent strip
    accent = resolve_color("accent", brand)
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(MARGIN_LEFT + 3.0),
        Inches(4.2),
        Inches(CONTENT_WIDTH - 6.0),
        Inches(0.06),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = accent
    shape.line.fill.background()


def _render_stat_callout(slide, node: SlideNode, brand: BrandConfig):
    """Render a stat callout slide with big numbers."""
    text_color = _text_color_for_bg(node.background, brand)
    muted = _muted_color_for_bg(node.background, brand)
    accent = resolve_color("accent", brand)

    # Heading
    if node.heading:
        _add_textbox(
            slide,
            TITLE_LEFT,
            TITLE_TOP,
            TITLE_WIDTH,
            TITLE_HEIGHT,
            node.heading,
            font_size=FONT_HEADING,
            bold=True,
            color=text_color,
            font_name=brand.header_font,
        )

    if not node.stats:
        return

    stat_count = len(node.stats)
    stat_width = CONTENT_WIDTH / stat_count
    stat_top = CONTENT_TOP + 0.5

    for i, stat in enumerate(node.stats):
        x = MARGIN_LEFT + i * stat_width

        # Value
        _add_textbox(
            slide,
            x,
            stat_top,
            stat_width,
            1.0,
            stat.value,
            font_size=FONT_STAT_VALUE,
            bold=True,
            color=accent,
            alignment=PP_ALIGN.CENTER,
            font_name=brand.header_font,
        )

        # Label
        _add_textbox(
            slide,
            x,
            stat_top + 1.2,
            stat_width,
            0.5,
            stat.label,
            font_size=FONT_STAT_LABEL,
            bold=True,
            color=text_color,
            alignment=PP_ALIGN.CENTER,
            font_name=brand.body_font,
        )

        # Description
        if stat.description:
            _add_textbox(
                slide,
                x,
                stat_top + 1.8,
                stat_width,
                0.5,
                stat.description,
                font_size=FONT_STAT_DESC,
                color=muted,
                alignment=PP_ALIGN.CENTER,
                font_name=brand.body_font,
            )


def _render_bullet_points(slide, node: SlideNode, brand: BrandConfig):
    """Render a bullet points slide."""
    text_color = _text_color_for_bg(node.background, brand)

    # Heading
    if node.heading:
        _add_textbox(
            slide,
            TITLE_LEFT,
            TITLE_TOP,
            TITLE_WIDTH,
            TITLE_HEIGHT,
            node.heading,
            font_size=FONT_HEADING,
            bold=True,
            color=text_color,
            font_name=brand.header_font,
        )

    if node.bullets:
        if node.layout == "icon_rows":
            _render_icon_rows(slide, node, brand)
        else:
            _add_bullet_list(
                slide,
                MARGIN_LEFT,
                CONTENT_TOP,
                CONTENT_WIDTH,
                4.5,
                node.bullets,
                font_size=FONT_BODY,
                color=text_color,
                font_name=brand.body_font,
            )


def _render_icon_rows(slide, node: SlideNode, brand: BrandConfig):
    """Render bullet points as icon rows layout."""
    text_color = _text_color_for_bg(node.background, brand)
    accent = resolve_color("accent", brand)
    row_height = 0.7
    start_top = CONTENT_TOP + 0.2

    for i, bullet in enumerate(node.bullets):
        y = start_top + i * row_height

        # Icon circle
        shape = slide.shapes.add_shape(
            9,  # MSO_SHAPE.OVAL
            Inches(MARGIN_LEFT),
            Inches(y),
            Inches(0.45),
            Inches(0.45),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = accent
        shape.line.fill.background()

        # Icon label inside circle
        if bullet.icon:
            tf = shape.text_frame
            tf.paragraphs[0].text = bullet.icon[0].upper()
            tf.paragraphs[0].font.size = Pt(14)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.word_wrap = False

        # Text
        _add_textbox(
            slide,
            MARGIN_LEFT + 0.7,
            y + 0.05,
            CONTENT_WIDTH - 0.7,
            0.4,
            bullet.text,
            font_size=FONT_BODY,
            color=text_color,
            font_name=brand.body_font,
        )


def _render_two_column(slide, node: SlideNode, brand: BrandConfig):
    """Render a two-column slide."""
    text_color = _text_color_for_bg(node.background, brand)

    # Heading
    if node.heading:
        _add_textbox(
            slide,
            TITLE_LEFT,
            TITLE_TOP,
            TITLE_WIDTH,
            TITLE_HEIGHT,
            node.heading,
            font_size=FONT_HEADING,
            bold=True,
            color=text_color,
            font_name=brand.header_font,
        )

    col_width = (CONTENT_WIDTH - COLUMN_GAP) / 2

    for i, col in enumerate(node.columns[:2]):
        x = MARGIN_LEFT + i * (col_width + COLUMN_GAP)
        y = CONTENT_TOP

        # Column title
        if col.title:
            _add_textbox(
                slide,
                x,
                y,
                col_width,
                0.5,
                col.title,
                font_size=FONT_BODY + 2,
                bold=True,
                color=text_color,
                font_name=brand.header_font,
            )
            y += 0.6

        # Column bullets
        if col.bullets:
            _add_bullet_list(
                slide,
                x,
                y,
                col_width,
                4.0,
                col.bullets,
                font_size=FONT_BODY,
                color=text_color,
                font_name=brand.body_font,
            )

    # Divider line
    if len(node.columns) >= 2:
        center_x = MARGIN_LEFT + col_width + COLUMN_GAP / 2
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(center_x - 0.01),
            Inches(CONTENT_TOP),
            Inches(0.02),
            Inches(4.5),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        shape.line.fill.background()


def _render_comparison(slide, node: SlideNode, brand: BrandConfig):
    """Render a comparison table slide."""
    text_color = _text_color_for_bg(node.background, brand)

    # Heading
    if node.heading:
        _add_textbox(
            slide,
            TITLE_LEFT,
            TITLE_TOP,
            TITLE_WIDTH,
            TITLE_HEIGHT,
            node.heading,
            font_size=FONT_HEADING,
            bold=True,
            color=text_color,
            font_name=brand.header_font,
        )

    if not node.compare:
        return

    headers = node.compare.headers
    rows = node.compare.rows
    col_count = len(headers) if headers else (len(rows[0]) if rows else 0)
    if col_count == 0:
        return

    row_count = len(rows) + (1 if headers else 0)
    table_top = CONTENT_TOP + 0.2
    table_height = min(row_count * 0.6, 4.5)

    table_shape = slide.shapes.add_table(
        row_count,
        col_count,
        Inches(MARGIN_LEFT),
        Inches(table_top),
        Inches(CONTENT_WIDTH),
        Inches(table_height),
    )
    table = table_shape.table

    # Header row
    if headers:
        for j, hdr in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = hdr
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(FONT_BODY)
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.font.name = brand.header_font
            cell.fill.solid()
            cell.fill.fore_color.rgb = resolve_color("primary", brand)

    # Data rows
    row_offset = 1 if headers else 0
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            if j >= col_count:
                break
            cell = table.cell(i + row_offset, j)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(FONT_BODY - 1)
            p.font.name = brand.body_font
            # Alternating row colors
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)


def _render_timeline(slide, node: SlideNode, brand: BrandConfig):
    """Render a timeline slide with horizontal steps."""
    text_color = _text_color_for_bg(node.background, brand)
    accent = resolve_color("accent", brand)
    muted = _muted_color_for_bg(node.background, brand)

    # Heading
    if node.heading:
        _add_textbox(
            slide,
            TITLE_LEFT,
            TITLE_TOP,
            TITLE_WIDTH,
            TITLE_HEIGHT,
            node.heading,
            font_size=FONT_HEADING,
            bold=True,
            color=text_color,
            font_name=brand.header_font,
        )

    if not node.timeline:
        return

    step_count = len(node.timeline)
    step_width = CONTENT_WIDTH / step_count
    line_y = CONTENT_TOP + 1.0

    # Connecting line
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(MARGIN_LEFT + 0.3),
        Inches(line_y + 0.18),
        Inches(CONTENT_WIDTH - 0.6),
        Inches(0.04),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = resolve_color("secondary", brand)
    shape.line.fill.background()

    for i, step in enumerate(node.timeline):
        x = MARGIN_LEFT + i * step_width + step_width / 2

        # Circle marker
        circle = slide.shapes.add_shape(
            9,  # MSO_SHAPE.OVAL
            Inches(x - 0.2),
            Inches(line_y),
            Inches(0.4),
            Inches(0.4),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent
        circle.line.fill.background()

        # Step number
        tf = circle.text_frame
        tf.paragraphs[0].text = str(i + 1)
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Time label
        _add_textbox(
            slide,
            x - step_width / 2,
            line_y - 0.6,
            step_width,
            0.4,
            step.time,
            font_size=FONT_CAPTION,
            bold=True,
            color=accent,
            alignment=PP_ALIGN.CENTER,
            font_name=brand.body_font,
        )

        # Title
        _add_textbox(
            slide,
            x - step_width / 2,
            line_y + 0.6,
            step_width,
            0.4,
            step.title,
            font_size=FONT_BODY - 1,
            bold=True,
            color=text_color,
            alignment=PP_ALIGN.CENTER,
            font_name=brand.body_font,
        )

        # Description
        if step.description:
            _add_textbox(
                slide,
                x - step_width / 2,
                line_y + 1.1,
                step_width,
                0.5,
                step.description,
                font_size=FONT_CAPTION,
                color=muted,
                alignment=PP_ALIGN.CENTER,
                font_name=brand.body_font,
            )


def _render_image_text(slide, node: SlideNode, brand: BrandConfig):
    """Render an image + text slide."""
    text_color = _text_color_for_bg(node.background, brand)

    # Heading
    if node.heading:
        _add_textbox(
            slide,
            TITLE_LEFT,
            TITLE_TOP,
            TITLE_WIDTH,
            TITLE_HEIGHT,
            node.heading,
            font_size=FONT_HEADING,
            bold=True,
            color=text_color,
            font_name=brand.header_font,
        )

    # Image placeholder (left side)
    if node.image:
        try:
            slide.shapes.add_picture(
                node.image,
                Inches(MARGIN_LEFT),
                Inches(CONTENT_TOP),
                Inches(CONTENT_WIDTH / 2 - COLUMN_GAP / 2),
                Inches(4.5),
            )
        except Exception:
            logger.warning("Image not found: %s, using placeholder", node.image)
            _add_textbox(
                slide,
                MARGIN_LEFT,
                CONTENT_TOP + 1.5,
                CONTENT_WIDTH / 2 - COLUMN_GAP / 2,
                1.0,
                f"[Image: {node.image}]",
                font_size=FONT_CAPTION,
                color=_muted_color_for_bg(node.background, brand),
                alignment=PP_ALIGN.CENTER,
            )

    # Text (right side)
    text_x = MARGIN_LEFT + CONTENT_WIDTH / 2 + COLUMN_GAP / 2
    text_width = CONTENT_WIDTH / 2 - COLUMN_GAP / 2

    if node.bullets:
        _add_bullet_list(
            slide,
            text_x,
            CONTENT_TOP,
            text_width,
            4.5,
            node.bullets,
            font_size=FONT_BODY,
            color=text_color,
            font_name=brand.body_font,
        )
    elif node.body:
        _add_textbox(
            slide,
            text_x,
            CONTENT_TOP,
            text_width,
            4.5,
            node.body,
            font_size=FONT_BODY,
            color=text_color,
            font_name=brand.body_font,
        )


def _render_quote(slide, node: SlideNode, brand: BrandConfig):
    """Render a quote slide."""
    text_color = _text_color_for_bg(node.background, brand)
    accent = resolve_color("accent", brand)
    muted = _muted_color_for_bg(node.background, brand)

    # Large quotation mark
    _add_textbox(
        slide,
        MARGIN_LEFT + 1.0,
        1.5,
        1.5,
        1.5,
        "\u201c",
        font_size=96,
        color=accent,
        font_name=brand.header_font,
    )

    # Quote text
    if node.heading:
        _add_textbox(
            slide,
            MARGIN_LEFT + 1.5,
            2.5,
            CONTENT_WIDTH - 3.0,
            2.5,
            node.heading,
            font_size=24,
            color=text_color,
            alignment=PP_ALIGN.CENTER,
            font_name=brand.body_font,
        )

    # Attribution
    if node.subheading:
        _add_textbox(
            slide,
            MARGIN_LEFT + 1.5,
            5.0,
            CONTENT_WIDTH - 3.0,
            0.5,
            f"\u2014 {node.subheading}",
            font_size=FONT_BODY,
            color=muted,
            alignment=PP_ALIGN.RIGHT,
            font_name=brand.body_font,
        )


def _render_closing(slide, node: SlideNode, brand: BrandConfig):
    """Render a closing slide."""
    text_color = _text_color_for_bg(node.background, brand)
    muted = _muted_color_for_bg(node.background, brand)

    if node.heading:
        _add_textbox(
            slide,
            MARGIN_LEFT,
            2.5,
            CONTENT_WIDTH,
            1.5,
            node.heading,
            font_size=FONT_TITLE,
            bold=True,
            color=text_color,
            alignment=PP_ALIGN.CENTER,
            font_name=brand.header_font,
        )

    if node.subheading:
        _add_textbox(
            slide,
            MARGIN_LEFT,
            4.2,
            CONTENT_WIDTH,
            0.8,
            node.subheading,
            font_size=FONT_SUBTITLE,
            color=muted,
            alignment=PP_ALIGN.CENTER,
            font_name=brand.body_font,
        )


def _render_freeform(slide, node: SlideNode, brand: BrandConfig):
    """Render a freeform slide -- best effort based on available content."""
    text_color = _text_color_for_bg(node.background, brand)

    if node.heading:
        _add_textbox(
            slide,
            TITLE_LEFT,
            TITLE_TOP,
            TITLE_WIDTH,
            TITLE_HEIGHT,
            node.heading,
            font_size=FONT_HEADING,
            bold=True,
            color=text_color,
            font_name=brand.header_font,
        )

    if node.bullets:
        _add_bullet_list(
            slide,
            MARGIN_LEFT,
            CONTENT_TOP,
            CONTENT_WIDTH,
            4.5,
            node.bullets,
            font_size=FONT_BODY,
            color=text_color,
            font_name=brand.body_font,
        )
    elif node.body:
        _add_textbox(
            slide,
            MARGIN_LEFT,
            CONTENT_TOP,
            CONTENT_WIDTH,
            4.5,
            node.body,
            font_size=FONT_BODY,
            color=text_color,
            font_name=brand.body_font,
        )


# ── Consulting Metadata Renderers ─────────────────────────────────


def _render_source_line(slide, node: SlideNode, brand: BrandConfig):
    """Render the source attribution line at the bottom of a slide."""
    if not node.source:
        return
    muted = _muted_color_for_bg(node.background, brand)
    _add_textbox(
        slide,
        MARGIN_LEFT,
        SOURCE_TOP,
        CONTENT_WIDTH - PAGE_NUM_WIDTH - 0.2,
        SOURCE_HEIGHT,
        node.source,
        font_size=SOURCE_FONT,
        color=muted,
        alignment=PP_ALIGN.LEFT,
        font_name=brand.body_font,
    )


def _render_exhibit_label(slide, node: SlideNode, brand: BrandConfig):
    """Render the exhibit label above the body content."""
    if not node.exhibit_label:
        return
    muted = _muted_color_for_bg(node.background, brand)
    _add_textbox(
        slide,
        MARGIN_LEFT,
        CONTENT_TOP - 0.3,
        CONTENT_WIDTH,
        0.25,
        node.exhibit_label,
        font_size=FONT_CAPTION,
        bold=True,
        color=muted,
        alignment=PP_ALIGN.LEFT,
        font_name=brand.body_font,
    )


def _render_footnotes(slide, node: SlideNode, brand: BrandConfig):
    """Render footnotes between body and source line."""
    if not node.footnotes:
        return
    muted = _muted_color_for_bg(node.background, brand)
    footnote_text = "  ".join(f"{i + 1}. {fn}" for i, fn in enumerate(node.footnotes))
    _add_textbox(
        slide,
        MARGIN_LEFT,
        SOURCE_TOP - 0.3,
        CONTENT_WIDTH,
        0.25,
        footnote_text,
        font_size=SOURCE_FONT,
        color=muted,
        alignment=PP_ALIGN.LEFT,
        font_name=brand.body_font,
    )


def _render_page_number(slide, page_num: int, node: SlideNode, brand: BrandConfig):
    """Render page number at bottom-right of slide."""
    muted = _muted_color_for_bg(node.background, brand)
    _add_textbox(
        slide,
        13.333 - MARGIN_RIGHT - PAGE_NUM_WIDTH,
        SOURCE_TOP,
        PAGE_NUM_WIDTH,
        SOURCE_HEIGHT,
        str(page_num),
        font_size=SOURCE_FONT,
        color=muted,
        alignment=PP_ALIGN.RIGHT,
        font_name=brand.body_font,
    )


# ── New Slide Type Renderers ─────────────────────────────────────


def _render_exec_summary(slide, node: SlideNode, brand: BrandConfig):
    """Render an executive summary slide with key messages."""
    text_color = _text_color_for_bg(node.background, brand)

    # Heading
    if node.heading:
        _add_textbox(
            slide,
            TITLE_LEFT,
            TITLE_TOP,
            TITLE_WIDTH,
            TITLE_HEIGHT,
            node.heading,
            font_size=FONT_HEADING,
            bold=True,
            color=text_color,
            font_name=brand.header_font,
        )

    # Executive summary bullets (key messages)
    if node.bullets:
        _add_bullet_list(
            slide,
            MARGIN_LEFT,
            CONTENT_TOP,
            CONTENT_WIDTH,
            4.5,
            node.bullets,
            font_size=FONT_BODY,
            color=text_color,
            font_name=brand.body_font,
        )


def _render_next_steps(slide, node: SlideNode, brand: BrandConfig):
    """Render a next-steps slide as an action item table."""
    text_color = _text_color_for_bg(node.background, brand)

    # Heading
    if node.heading:
        _add_textbox(
            slide,
            TITLE_LEFT,
            TITLE_TOP,
            TITLE_WIDTH,
            TITLE_HEIGHT,
            node.heading,
            font_size=FONT_HEADING,
            bold=True,
            color=text_color,
            font_name=brand.header_font,
        )

    if not node.next_steps:
        # Fall back to bullets if no @action directives
        if node.bullets:
            _add_bullet_list(
                slide,
                MARGIN_LEFT,
                CONTENT_TOP,
                CONTENT_WIDTH,
                4.5,
                node.bullets,
                font_size=FONT_BODY,
                color=text_color,
                font_name=brand.body_font,
            )
        return

    # Render as table: Action | Owner | Timeline
    headers = ["Action", "Owner", "Timeline"]
    row_count = len(node.next_steps) + 1  # +1 for header
    col_count = 3
    table_top = CONTENT_TOP + 0.2
    table_height = min(row_count * 0.6, 4.5)

    table_shape = slide.shapes.add_table(
        row_count,
        col_count,
        Inches(MARGIN_LEFT),
        Inches(table_top),
        Inches(CONTENT_WIDTH),
        Inches(table_height),
    )
    table = table_shape.table

    # Header row
    for j, hdr in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = hdr
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(FONT_BODY)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.name = brand.header_font
        cell.fill.solid()
        cell.fill.fore_color.rgb = resolve_color("primary", brand)

    # Data rows
    for i, ns in enumerate(node.next_steps):
        values = [ns.action, ns.owner or "", ns.timeline or ""]
        for j, val in enumerate(values):
            cell = table.cell(i + 1, j)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(FONT_BODY - 1)
            p.font.name = brand.body_font
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)


# ── Speaker Notes ─────────────────────────────────────────────────


def _add_speaker_notes(slide, notes: str):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes


# ── Dispatch Table ────────────────────────────────────────────────

_RENDERERS = {
    SlideType.TITLE: _render_title,
    SlideType.SECTION_DIVIDER: _render_section_divider,
    SlideType.STAT_CALLOUT: _render_stat_callout,
    SlideType.BULLET_POINTS: _render_bullet_points,
    SlideType.TWO_COLUMN: _render_two_column,
    SlideType.COMPARISON: _render_comparison,
    SlideType.TIMELINE: _render_timeline,
    SlideType.IMAGE_TEXT: _render_image_text,
    SlideType.QUOTE: _render_quote,
    SlideType.CLOSING: _render_closing,
    SlideType.EXEC_SUMMARY: _render_exec_summary,
    SlideType.NEXT_STEPS: _render_next_steps,
    SlideType.FREEFORM: _render_freeform,
}


# ── Public API ────────────────────────────────────────────────────


def render(
    presentation: PresentationNode,
    output_dir: Path,
    template_path: Optional[str] = None,
) -> Path:
    """Render a PresentationNode to a .pptx file.

    Args:
        presentation: Parsed presentation from DSL.
        output_dir: Directory to write output file.
        template_path: Optional .pptx template (not yet implemented).

    Returns:
        Path to the generated .pptx file.
    """
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    brand = presentation.meta.brand
    blank_layout = prs.slide_layouts[6]  # blank layout

    for page_num, node in enumerate(presentation.slides, start=1):
        slide = prs.slides.add_slide(blank_layout)

        # Background
        _apply_background(slide, node.background, brand)

        # Content via dispatch
        renderer_fn = _RENDERERS.get(node.slide_type, _render_freeform)
        renderer_fn(slide, node, brand)

        # Consulting metadata (exhibit label, footnotes, source, page number)
        _render_exhibit_label(slide, node, brand)
        _render_footnotes(slide, node, brand)
        _render_source_line(slide, node, brand)
        _render_page_number(slide, page_num, node, brand)

        # Speaker notes
        if node.speaker_notes:
            _add_speaker_notes(slide, node.speaker_notes)

    # Save
    safe_title = "".join(
        c if c.isalnum() or c in " -_" else "_" for c in presentation.meta.title
    ).strip()[:80]
    filename = f"{safe_title}.pptx"
    output_path = output_dir / filename
    prs.save(str(output_path))

    logger.info("Rendered %d slides to %s", len(presentation.slides), output_path)
    return output_path
