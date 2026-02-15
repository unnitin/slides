"""
skills/template_analyze.py — Introspect a .pptx template for available layouts.

Uses python-pptx to extract layout names, placeholder types, and dimensions
from an existing .pptx template file.
"""

from pptx import Presentation


def analyze_template(template_path: str) -> dict:
    """Analyze a .pptx template and return its structure.

    Args:
        template_path: Path to a .pptx template file.

    Returns:
        Dict with slide_width, slide_height, and layouts list.
        Each layout includes name, index, and placeholder details.
    """
    prs = Presentation(template_path)

    layouts = []
    for i, layout in enumerate(prs.slide_layouts):
        placeholders = []
        for ph in layout.placeholders:
            placeholders.append(
                {
                    "idx": ph.placeholder_format.idx,
                    "name": ph.name,
                    "type": str(ph.placeholder_format.type),
                    "left": _emu_to_inches(ph.left),
                    "top": _emu_to_inches(ph.top),
                    "width": _emu_to_inches(ph.width),
                    "height": _emu_to_inches(ph.height),
                }
            )
        layouts.append(
            {
                "index": i,
                "name": layout.name,
                "placeholder_count": len(placeholders),
                "placeholders": placeholders,
            }
        )

    return {
        "slide_width": _emu_to_inches(prs.slide_width),
        "slide_height": _emu_to_inches(prs.slide_height),
        "layout_count": len(layouts),
        "layouts": layouts,
    }


def _emu_to_inches(emu: int) -> float:
    """Convert EMU (English Metric Units) to inches."""
    return round(emu / 914400, 2)
